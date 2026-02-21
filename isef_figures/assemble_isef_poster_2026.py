"""
ISEF 2026 - NOD2 Project - POSTER ASSEMBLY
Author: Tanmay Vasudeva
Texas Virtual Academy at Hallsville

Creates a professional 36x48 ISEF poster following the winning plan:
- Title and one-sentence story
- 3 headline results (FEP-focused)
- Pipeline, MD dynamics, FEP charts
- PyMOL structure images
- Conclusions and future work
"""

import os
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("ERROR: python-pptx not installed. Run: pip install python-pptx")
    exit(1)

# ============== PATHS ==============
BASE_DIR = Path("C:/Users/vasud/nod2-screening-data/isef_figures")
FIGURES_DIR = BASE_DIR / "winner_figures"
PATHWAY_DIR = BASE_DIR / "pathway_diagram"
OUTPUT_DIR = BASE_DIR / "final_poster"
OUTPUT_DIR.mkdir(exist_ok=True)

# Template (36x48 inches at 96 DPI)
# Standard ISEF poster dimensions
POSTER_WIDTH = Inches(36)
POSTER_HEIGHT = Inches(48)

# ============== COLOR SCHEME ==============
NAVY = RGBColor(0x1a, 0x36, 0x5d)      # #1a365d - Headers
TEAL = RGBColor(0x31, 0x97, 0x95)      # #319795 - Accents
CORAL = RGBColor(0xed, 0x89, 0x36)     # #ed8936 - Key findings
WHITE = RGBColor(0xff, 0xff, 0xff)
DARK_TEXT = RGBColor(0x1a, 0x20, 0x2c) # #1a202c
LIGHT_BG = RGBColor(0xf7, 0xfa, 0xfc)  # #f7fafc - Section backgrounds


def create_poster():
    """Create the ISEF 2026 poster."""

    # Create new presentation with custom size
    prs = Presentation()
    prs.slide_width = POSTER_WIDTH
    prs.slide_height = POSTER_HEIGHT

    # Add blank slide
    blank_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_layout)

    # ================================================================
    # TITLE SECTION (Top 6 inches)
    # ================================================================

    # Title background
    title_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), POSTER_WIDTH, Inches(6)
    )
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = NAVY
    title_bg.line.fill.background()

    # Main title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(35), Inches(1.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Deep Learning Guided Discovery and FEP Validation\nof NOD2 Binders for the R702W Crohn's Variant"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Author info
    author_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(35), Inches(0.8))
    tf = author_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Tanmay Vasudeva  •  Texas Virtual Academy at Hallsville  •  ISEF 2026"
    p.font.size = Pt(36)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # One-sentence story box
    story_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(4.2), Inches(32), Inches(1.4)
    )
    story_box.fill.solid()
    story_box.fill.fore_color.rgb = CORAL
    story_box.line.fill.background()

    story_tf = story_box.text_frame
    story_tf.word_wrap = True
    story_p = story_tf.paragraphs[0]
    story_p.text = "I built a screening pipeline for NOD2 binders and used FEP to quantify how the R702W variant changes binding for specific ligands, with apo MD providing supporting context."
    story_p.font.size = Pt(28)
    story_p.font.bold = True
    story_p.font.color.rgb = WHITE
    story_p.alignment = PP_ALIGN.CENTER
    story_tf.paragraphs[0].space_before = Pt(12)

    # ================================================================
    # 3 HEADLINE CONTRIBUTIONS (Row below title, ~2 inches)
    # ================================================================

    headline_y = Inches(6.5)
    headline_height = Inches(1.8)
    headline_width = Inches(10.5)

    headlines = [
        ("1", "FEP: Febuxostat Mutation-Sensitive", "ΔΔG = +2.34 kcal/mol, 50× weaker (p<0.001)"),
        ("2", "FEP: Bufadienolide No Change", "ΔΔG = -0.44 kcal/mol (NS)"),
        ("3", "Apo MD: Reduced R702W Sampling", "WT: 12.6% vs R702W: 0% frames >5Å"),
    ]

    for i, (num, title, detail) in enumerate(headlines):
        x = Inches(1 + i * 11.5)
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, headline_y, headline_width, headline_height
        )
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_BG
        box.line.color.rgb = TEAL
        box.line.width = Pt(3)

        # Number circle
        num_box = slide.shapes.add_textbox(x + Inches(0.2), headline_y + Inches(0.2), Inches(0.6), Inches(0.6))
        num_tf = num_box.text_frame
        num_p = num_tf.paragraphs[0]
        num_p.text = num
        num_p.font.size = Pt(32)
        num_p.font.bold = True
        num_p.font.color.rgb = CORAL

        # Title
        title_tb = slide.shapes.add_textbox(x + Inches(0.9), headline_y + Inches(0.15), Inches(9), Inches(0.7))
        title_tf = title_tb.text_frame
        title_p = title_tf.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(24)
        title_p.font.bold = True
        title_p.font.color.rgb = NAVY

        # Detail
        detail_tb = slide.shapes.add_textbox(x + Inches(0.9), headline_y + Inches(0.85), Inches(9.2), Inches(0.8))
        detail_tf = detail_tb.text_frame
        detail_p = detail_tf.paragraphs[0]
        detail_p.text = detail
        detail_p.font.size = Pt(20)
        detail_p.font.color.rgb = DARK_TEXT

    # ================================================================
    # MAIN CONTENT AREA - LEFT COLUMN (Pipeline + Methods)
    # ================================================================

    left_x = Inches(0.5)
    col_width = Inches(11)
    content_y = Inches(9)

    # Section: Introduction & Hypothesis
    add_section_header(slide, "Introduction", left_x, content_y, col_width)

    intro_text = """• NOD2 mutations cause 15-25% of Crohn's disease
• R702W is the most common mutation (~10% of patients)
• Literature describes R702W as "destabilizing" NOD2
• Hypothesis: Test binding differences at molecular level"""

    add_text_box(slide, intro_text, left_x, content_y + Inches(0.7), col_width, Inches(2.5))

    # Pipeline figure
    pipeline_fig = FIGURES_DIR / "figure1_pipeline_professional.png"
    if pipeline_fig.exists():
        slide.shapes.add_picture(str(pipeline_fig), left_x, content_y + Inches(3.5), width=col_width)

    # Section: Methods (Below pipeline)
    methods_y = content_y + Inches(7)
    add_section_header(slide, "Methods", left_x, methods_y, col_width)

    methods_text = """• 9,566 compounds (2,152 FDA + 7,414 natural products)
• GNINA CNN-docking on AlphaFold2 NOD2 LRR
• NOD2-Scout ML: AUC 0.85-0.93 (scaffold-split CV)
• MD: OpenMM, AMBER14 + OpenFF 2.1.0, 520 ns total
• FEP: Gold-standard thermodynamic binding validation"""

    add_text_box(slide, methods_text, left_x, methods_y + Inches(0.7), col_width, Inches(3))

    # ================================================================
    # MAIN CONTENT AREA - CENTER COLUMN (FEP Results - HEADLINE)
    # ================================================================

    center_x = Inches(12.5)

    # Section: FEP Binding Results (KEY FINDING)
    add_section_header(slide, "FEP Binding Results (Gold Standard)", center_x, content_y, col_width)

    # FEP figure
    fep_fig = FIGURES_DIR / "figure4_fep_comparison.png"
    if fep_fig.exists():
        slide.shapes.add_picture(str(fep_fig), center_x, content_y + Inches(0.8), width=Inches(10))

    # FEP data table
    fep_table_y = content_y + Inches(6.5)
    add_section_header(slide, "FEP Quantitative Results", center_x, fep_table_y, col_width)

    fep_table_text = """Compound         ΔG(WT)      ΔG(R702W)    ΔΔG         Significance
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
Febuxostat      -10.36       -8.02      +2.34 kcal/mol   ***  (8σ)
Bufadienolide   -15.22      -15.66      -0.44 kcal/mol    NS  (1.8σ)

Key: *** = p<0.001 (highly significant), NS = not significant"""

    add_text_box(slide, fep_table_text, center_x, fep_table_y + Inches(0.7), col_width, Inches(2.5), font_size=18)

    # Apo MD section
    md_y = fep_table_y + Inches(4)
    add_section_header(slide, "Apo MD: Supporting Context", center_x, md_y, col_width)

    # Dynamics figures (side by side)
    rmsd_fig = FIGURES_DIR / "figure2_rmsd_clean.png"
    het_fig = FIGURES_DIR / "figure3_heterogeneity_clean.png"

    if rmsd_fig.exists():
        slide.shapes.add_picture(str(rmsd_fig), center_x, md_y + Inches(0.8), width=Inches(5))
    if het_fig.exists():
        slide.shapes.add_picture(str(het_fig), center_x + Inches(5.5), md_y + Inches(0.8), width=Inches(5))

    # ================================================================
    # MAIN CONTENT AREA - RIGHT COLUMN (Structure + Conclusions)
    # ================================================================

    right_x = Inches(24.5)

    # Section: 3D Structure Visualization
    add_section_header(slide, "NOD2 Structure & Binding Site", right_x, content_y, col_width)

    # Placeholder for PyMOL images
    struct_note = slide.shapes.add_textbox(right_x, content_y + Inches(0.8), col_width, Inches(5))
    tf = struct_note.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "[Insert PyMOL structure images here]\n\n• Run pathway_diagram_10K.pml in PyMOL\n• Insert img1-img6 from pathway_diagram folder\n\nDomain colors:\n• Gray: CARD domains\n• Blue: NBD\n• Cyan: HD1\n• Orange: HD2 (R702W location)\n• Green: LRR (binding domain)"
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_TEXT

    # MD Validation section
    validation_y = content_y + Inches(7)
    add_section_header(slide, "MD Binding Validation", right_x, validation_y, col_width)

    occupancy_fig = FIGURES_DIR / "figure8_pocket_occupancy.png"
    if occupancy_fig.exists():
        slide.shapes.add_picture(str(occupancy_fig), right_x, validation_y + Inches(0.8), width=Inches(10))

    # ================================================================
    # BOTTOM SECTION - Conclusions & Future Work
    # ================================================================

    bottom_y = Inches(40)

    # Conclusions
    add_section_header(slide, "Conclusions", left_x, bottom_y, col_width)

    conclusions_text = """1. Built deep learning pipeline screening 9,566 compounds against NOD2 LRR
2. FEP quantifies mutation-sensitive binding:
   • Febuxostat: ΔΔG = +2.34 kcal/mol (50× weaker in R702W, p<0.001)
   • Bufadienolide: ΔΔG = -0.44 kcal/mol (no significant change)
3. Apo MD shows reduced conformational sampling in R702W (0% vs 12.6% >5Å)
4. Identified key binding residues: GLU1008, ARG1037, ASP1011"""

    add_text_box(slide, conclusions_text, left_x, bottom_y + Inches(0.7), Inches(16), Inches(3.5))

    # Limitations
    limitations_y = bottom_y
    add_section_header(slide, "Limitations", center_x, limitations_y, Inches(8))

    limitations_text = """• FEP on 2 ligands only (compute-intensive)
• Apo MD only (no drug comparisons)
• Computational predictions require
  experimental validation
• Bufadienolide statistical significance
  not reached (1.8σ)"""

    add_text_box(slide, limitations_text, center_x, limitations_y + Inches(0.7), Inches(8), Inches(3))

    # Future Work
    future_y = bottom_y
    add_section_header(slide, "Future Work", right_x, future_y, col_width)

    future_text = """• Experimental validation: SPR, ITC binding assays
• Test additional NOD2 mutations (G908R, L1007fs)
• HDX-MS conformational dynamics studies
• Cell-based NF-κB reporter assays
• Analog screening for improved drug-likeness"""

    add_text_box(slide, future_text, right_x, future_y + Inches(0.7), col_width, Inches(3))

    # ================================================================
    # FOOTER
    # ================================================================

    footer_y = Inches(45)

    # References/Acknowledgments
    footer_text = "Acknowledgments: TACC Frontera allocation, AlphaFold2 structure (UniProt Q9HC29), OpenMM/AMBER14 force fields"
    footer_box = slide.shapes.add_textbox(Inches(0.5), footer_y, Inches(25), Inches(0.5))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = footer_text
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_TEXT

    # QR code placeholder
    qr_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(32), footer_y - Inches(1), Inches(3), Inches(3)
    )
    qr_box.fill.solid()
    qr_box.fill.fore_color.rgb = LIGHT_BG
    qr_box.line.color.rgb = NAVY

    qr_label = slide.shapes.add_textbox(Inches(32), footer_y + Inches(2.2), Inches(3), Inches(0.5))
    tf = qr_label.text_frame
    p = tf.paragraphs[0]
    p.text = "GitHub\nRepository"
    p.font.size = Pt(14)
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER

    # ================================================================
    # SAVE
    # ================================================================

    output_path = OUTPUT_DIR / "ISEF_2026_NOD2_Poster.pptx"
    prs.save(str(output_path))
    print(f"Poster saved to: {output_path}")

    return output_path


def add_section_header(slide, text, x, y, width):
    """Add a section header with styling."""
    header = slide.shapes.add_textbox(x, y, width, Inches(0.5))
    tf = header.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = NAVY


def add_text_box(slide, text, x, y, width, height, font_size=20):
    """Add a text box with the given content."""
    box = slide.shapes.add_textbox(x, y, width, height)
    tf = box.text_frame
    tf.word_wrap = True

    # Split by lines and add paragraphs
    lines = text.strip().split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(6)


if __name__ == "__main__":
    print("=" * 60)
    print("ISEF 2026 - NOD2 Project - POSTER ASSEMBLY")
    print("=" * 60)

    # Create the poster
    output_path = create_poster()

    print()
    print("=" * 60)
    print("POSTER CREATED SUCCESSFULLY!")
    print("=" * 60)
    print()
    print(f"Output: {output_path}")
    print()
    print("Next steps:")
    print("1. Open the PowerPoint file")
    print("2. Run PyMOL script: pathway_diagram_10K.pml")
    print("3. Insert PyMOL images into structure section")
    print("4. Review and adjust layout as needed")
    print("5. Export to PDF for printing")
