"""
ISEF 2026 Poster Assembly
Author: Tanmay Vasudeva
Texas Virtual Academy at Hallsville

Assembles all figures into 36x48 poster using PowerPoint.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.util import Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os

# Paths
TEMPLATE = "C:/Users/vasud/Downloads/PosterPresentations.com-36x48-Template-Newfield.pptx"
FIGURES_DIR = "C:/Users/vasud/nod2-screening-data/isef_figures/v4_isef_style/"
OUTPUT = "C:/Users/vasud/nod2-screening-data/isef_figures/ISEF_2026_Poster_NOD2.pptx"

# Colors
NAVY = RGBColor(26, 54, 93)      # #1a365d
TEAL = RGBColor(49, 151, 149)    # #319795
WHITE = RGBColor(255, 255, 255)

def create_poster():
    """Create the ISEF poster from scratch with 36x48 dimensions."""

    # Create new presentation with poster dimensions
    prs = Presentation()
    prs.slide_width = Inches(36)
    prs.slide_height = Inches(48)

    # Add blank slide
    blank_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_layout)

    # ==================== HEADER ====================
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(35), Inches(2))
    tf = title_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Deep Learning Guided Discovery and FEP Validation of NOD2 Binders for the R702W Crohn's Variant"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER

    # Author
    author_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.6), Inches(35), Inches(0.8))
    tf = author_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Tanmay Vasudeva"
    p.font.size = Pt(48)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # School
    school_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.3), Inches(35), Inches(0.6))
    tf = school_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Texas Virtual Academy at Hallsville"
    p.font.size = Pt(36)
    p.font.italic = True
    p.alignment = PP_ALIGN.CENTER

    # ==================== ONE-SENTENCE STORY ====================
    story_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Inches(1), Inches(4.2), Inches(34), Inches(1.2))
    story_box.fill.solid()
    story_box.fill.fore_color.rgb = RGBColor(240, 240, 240)
    story_box.line.color.rgb = NAVY

    story_text = slide.shapes.add_textbox(Inches(1.2), Inches(4.35), Inches(33.6), Inches(1))
    tf = story_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "I built a screening pipeline for NOD2 binders and used FEP to quantify how the R702W variant changes binding for specific ligands, with apo MD providing supporting context."
    p.font.size = Pt(28)
    p.font.italic = True
    p.alignment = PP_ALIGN.CENTER

    # ==================== COLUMN 1: INTRODUCTION & PIPELINE ====================
    col1_x = Inches(0.5)
    col_width = Inches(11)

    # Introduction Header
    add_section_header(slide, "Introduction", col1_x, Inches(5.8), col_width)

    intro_text = """• NOD2 mutations cause 15-25% of Crohn's disease cases
• R702W is the most common mutation (~10% of patients)
• Hypothesis: Test mutation effects at molecular level with MD and FEP"""
    add_text_box(slide, intro_text, col1_x, Inches(6.6), col_width, Inches(2.2), size=24)

    # Pipeline Header
    add_section_header(slide, "Computational Pipeline", col1_x, Inches(9), col_width)

    # Pipeline Figure
    pipeline_img = FIGURES_DIR + "fig_pipeline.png"
    if os.path.exists(pipeline_img):
        slide.shapes.add_picture(pipeline_img, col1_x, Inches(9.8), width=col_width)

    # Pipeline details
    pipeline_text = """• 9,566 compounds screened (FDA + natural products)
• GNINA CNN-based docking on AlphaFold NOD2 LRR
• NOD2-Scout ML model (AUC 0.85-0.93, scaffold-split CV)
• 520 ns total MD simulation (OpenMM, AMBER14)
• FEP binding validation (gold standard thermodynamics)"""
    add_text_box(slide, pipeline_text, col1_x, Inches(12.8), col_width, Inches(3), size=22)

    # ML Performance Figure
    ml_img = FIGURES_DIR + "fig_ml_performance.png"
    if os.path.exists(ml_img):
        slide.shapes.add_picture(ml_img, col1_x, Inches(16), width=Inches(10))

    # ==================== COLUMN 2: RESULTS ====================
    col2_x = Inches(12.3)

    # FEP Results Header
    add_section_header(slide, "FEP Binding Free Energy (Key Result)", col2_x, Inches(5.8), col_width)

    # FEP Figure
    fep_img = FIGURES_DIR + "fig_fep.png"
    if os.path.exists(fep_img):
        slide.shapes.add_picture(fep_img, Inches(12.8), Inches(6.6), width=Inches(10))

    # FEP Results text
    fep_text = """FEP Results (Gold Standard):
• Febuxostat: ΔΔG = +2.34 kcal/mol
  → 50× weaker binding in R702W (p<0.001)
• Bufadienolide: ΔΔG = -0.44 kcal/mol
  → No significant change (n.s.)

Key Finding: Ligand-dependent binding differences between WT and R702W"""
    add_text_box(slide, fep_text, col2_x, Inches(11.5), col_width, Inches(3.5), size=22)

    # MD Validation Header
    add_section_header(slide, "MD Binding Validation", col2_x, Inches(15.2), col_width)

    # Pocket Occupancy Figure
    pocket_img = FIGURES_DIR + "fig_pocket_occupancy.png"
    if os.path.exists(pocket_img):
        slide.shapes.add_picture(pocket_img, Inches(12.8), Inches(16), width=Inches(10))

    # ==================== COLUMN 3: DYNAMICS & CONCLUSIONS ====================
    col3_x = Inches(24.3)

    # Apo Dynamics Header
    add_section_header(slide, "Apo MD Dynamics (Supporting Context)", col3_x, Inches(5.8), col_width)

    # Dynamics Figure
    dynamics_img = FIGURES_DIR + "fig_dynamics.png"
    if os.path.exists(dynamics_img):
        slide.shapes.add_picture(dynamics_img, col3_x, Inches(6.6), width=col_width)

    dynamics_text = """• WT: 12.6% frames exceed 5Å RMSD threshold
• R702W: 0% frames exceed threshold
• R702W shows reduced conformational sampling
• Consistent with mutation-sensitive binding in FEP"""
    add_text_box(slide, dynamics_text, col3_x, Inches(11), col_width, Inches(2.5), size=22)

    # RMSD Trace
    trace_img = FIGURES_DIR + "fig_rmsd_trace.png"
    if os.path.exists(trace_img):
        slide.shapes.add_picture(trace_img, col3_x, Inches(13.8), width=col_width)

    # ==================== BOTTOM SECTIONS ====================

    # Conclusions
    add_section_header(slide, "Conclusions", col1_x, Inches(21), Inches(16))

    conclusions_text = """1. Built deep learning pipeline screening 9,566 compounds against NOD2 LRR domain

2. FEP quantifies mutation-sensitive binding:
   • Febuxostat: ΔΔG = +2.34 kcal/mol (50× weaker in R702W, p<0.001)
   • Bufadienolide: ΔΔG = -0.44 kcal/mol (no significant change)

3. Apo MD shows reduced conformational sampling in R702W vs WT (0% vs 12.6% frames >5Å)

4. Identified key binding residues: GLU1008, ARG1037, ASP1011"""
    add_text_box(slide, conclusions_text, col1_x, Inches(21.8), Inches(16), Inches(5), size=24)

    # Limitations
    add_section_header(slide, "Limitations", Inches(18), Inches(21), Inches(8))

    limitations_text = """• FEP performed on 2 ligands only
• Apo MD simulations only
• Computational predictions require
  experimental validation"""
    add_text_box(slide, limitations_text, Inches(18), Inches(21.8), Inches(8), Inches(2.5), size=20)

    # Future Work
    add_section_header(slide, "Future Work", Inches(27), Inches(21), Inches(8))

    future_text = """• SPR/ITC experimental validation
• Test additional mutations
  (G908R, L1007fs)
• HDX-MS experiments"""
    add_text_box(slide, future_text, Inches(27), Inches(21.8), Inches(8), Inches(2.5), size=20)

    # ==================== FOOTER ====================
    footer_text = "Generated with Claude Code | ISEF 2026"
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(46.5), Inches(35), Inches(0.5))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = footer_text
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(128, 128, 128)
    p.alignment = PP_ALIGN.CENTER

    # Save
    prs.save(OUTPUT)
    print(f"Poster saved to: {OUTPUT}")


def add_section_header(slide, text, x, y, width):
    """Add a section header with navy background."""
    # Background box
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, width, Inches(0.7))
    box.fill.solid()
    box.fill.fore_color.rgb = NAVY
    box.line.fill.background()

    # Text
    text_box = slide.shapes.add_textbox(x, y + Inches(0.1), width, Inches(0.6))
    tf = text_box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER


def add_text_box(slide, text, x, y, width, height, size=24):
    """Add a text box with bullet points."""
    text_box = slide.shapes.add_textbox(x, y, width, height)
    tf = text_box.text_frame
    tf.word_wrap = True

    lines = text.strip().split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.space_after = Pt(6)


if __name__ == "__main__":
    create_poster()
    print("\nPoster assembled successfully!")
    print("Open the file in PowerPoint to review and adjust as needed.")
