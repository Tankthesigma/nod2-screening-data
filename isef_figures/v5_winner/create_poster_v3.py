"""
ISEF 2026 Poster V3 - DENSE TEXT + VISUALS BALANCED
Like Rishab Jain's style: 60% text, 40% visuals
Body text: 32pt minimum
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(48)
prs.slide_height = Inches(36)
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Colors - Clean scientific
NAVY = RGBColor(0x1a, 0x36, 0x5d)
DARK_TEXT = RGBColor(0x1f, 0x2d, 0x3d)
ACCENT = RGBColor(0xc0, 0x39, 0x2b)  # Red for key findings
LIGHT_BG = RGBColor(0xf5, 0xf6, 0xf8)
WHITE = RGBColor(255, 255, 255)
MID_GRAY = RGBColor(0x64, 0x74, 0x8b)

def rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()

def text(slide, l, t, w, h, txt, size, color=DARK_TEXT, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = txt
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align

def header(slide, l, t, w, txt):
    rect(slide, l, t, w, Inches(0.55), NAVY)
    text(slide, l + Inches(0.1), t + Inches(0.07), w, Inches(0.45), txt.upper(), 32, WHITE, True)

# Background
rect(slide, Inches(0), Inches(0), Inches(48), Inches(36), LIGHT_BG)

# ============================================================
# TITLE BANNER
# ============================================================
rect(slide, Inches(0), Inches(0), Inches(48), Inches(3.2), NAVY)

text(slide, Inches(0.5), Inches(0.3), Inches(47), Inches(1.4),
     "Deep Learning Guided Discovery and FEP Validation of NOD2 Binders for the R702W Crohn's Variant",
     72, WHITE, True, PP_ALIGN.CENTER)

text(slide, Inches(0.5), Inches(1.6), Inches(47), Inches(0.7),
     "Tanmay Vasudeva  |  Texas Virtual Academy at Hallsville",
     36, WHITE, False, PP_ALIGN.CENTER)

# Key finding callout
rect(slide, Inches(8), Inches(2.3), Inches(32), Inches(0.7), ACCENT)
text(slide, Inches(8.2), Inches(2.35), Inches(31.6), Inches(0.6),
     "KEY FINDING: Febuxostat exhibits 50x weaker binding in R702W mutant (ddG = +2.34 kcal/mol, p<0.001)",
     28, WHITE, True, PP_ALIGN.CENTER)

# ============================================================
# LAYOUT: 3 columns
# ============================================================
top = Inches(3.5)
col1_l = Inches(0.5)
col1_w = Inches(14.5)
col2_l = Inches(15.5)
col2_w = Inches(16)
col3_l = Inches(32)
col3_w = Inches(15.5)

# ============================================================
# COLUMN 1: INTRO + BACKGROUND + METHODS
# ============================================================

header(slide, col1_l, top, col1_w, "Research Question")

q_text = """What molecular mechanisms explain how the Crohn's-associated NOD2 R702W variant alters drug binding affinity? This project tests whether the R702W substitution, located 80 angstroms from the LRR binding pocket, produces measurable differences in ligand binding through an allosteric mechanism. I developed a computational pipeline combining machine learning screening with rigorous free energy perturbation (FEP) calculations to quantify mutation-sensitive binding for FDA-approved compounds."""

text(slide, col1_l, top + Inches(0.6), col1_w, Inches(3.5), q_text, 28, DARK_TEXT)

header(slide, col1_l, top + Inches(4), col1_w, "Background")

bg_text = """NOD2 is a cytosolic pattern-recognition receptor that detects bacterial peptidoglycans and initiates NF-kB inflammatory signaling. Loss-of-function variants in NOD2 represent the strongest genetic risk factors for Crohn's disease, affecting 3+ million patients globally. The R702W substitution accounts for 15-25% of Crohn's cases, replacing a positively charged arginine with a bulky hydrophobic tryptophan in the HD2 domain. Despite its clinical significance, no targeted therapies exist for R702W carriers. Drug repurposing offers a faster path to therapy than de novo design, motivating this computational screening approach."""

text(slide, col1_l, top + Inches(4.6), col1_w, Inches(4.2), bg_text, 28, DARK_TEXT)

header(slide, col1_l, top + Inches(8.7), col1_w, "Computational Methods")

methods_text = """Structure Preparation: Full-length NOD2 models for wild-type and R702W were built using AlphaFold2 (Q9HC29). The LRR domain (residues 744-1040) was extracted for docking studies while retaining the HD2 domain context for mutation analysis.

Virtual Screening: 9,566 compounds from FDA-approved and natural product libraries were screened against the LRR binding pocket using GNINA, a CNN-based molecular docking program. A custom gradient-boosted classifier (NOD2-Scout) was trained to prioritize hits, achieving AUC of 0.85-0.93 on scaffold-split cross-validation.

Molecular Dynamics: Top-ranked complexes and apo structures were simulated using OpenMM with AMBER14 force field and TIP3P-FB water model at 310K. Total simulation time: 520 ns across all systems (20 ns production per replicate, n=3-4 per condition).

Free Energy Perturbation: Gold-standard FEP calculations quantified binding free energy differences between WT and R702W. Each ligand required 120 windows (80 complex-phase + 40 solvent-phase) with full thermodynamic cycle closure."""

text(slide, col1_l, top + Inches(9.3), col1_w, Inches(10), methods_text, 26, DARK_TEXT)

# Small pipeline figure
rect(slide, col1_l, top + Inches(19.2), col1_w, Inches(4.5), WHITE)
text(slide, col1_l, top + Inches(21), col1_w, Inches(0.8), "[Pipeline Flowchart]", 24, MID_GRAY, False, PP_ALIGN.CENTER)
text(slide, col1_l, top + Inches(23.4), col1_w, Inches(0.5), "Fig 1: Computational pipeline from AlphaFold to FEP validation", 20, DARK_TEXT, False, PP_ALIGN.CENTER)

header(slide, col1_l, top + Inches(24.2), col1_w, "Limitations")

limits = """FEP calculations were performed on two representative ligands due to computational cost (~120 GPU-hours per compound). All predictions require experimental validation through SPR or ITC binding assays. Apo MD simulations do not capture ligand-induced conformational changes."""

text(slide, col1_l, top + Inches(24.8), col1_w, Inches(2.5), limits, 26, DARK_TEXT)

# ============================================================
# COLUMN 2: RESULTS WITH INTERPRETATION
# ============================================================

header(slide, col2_l, top, col2_w, "FEP Binding Free Energy Results")

fep_intro = """Free energy perturbation (FEP) is the thermodynamic gold standard for computing binding affinity changes. I calculated the mutation effect as ddG = dG(R702W) - dG(WT), where positive values indicate weaker binding in the mutant."""

text(slide, col2_l, top + Inches(0.6), col2_w, Inches(1.8), fep_intro, 28, DARK_TEXT)

# FEP chart placeholder
rect(slide, col2_l, top + Inches(2.3), col2_w, Inches(5.5), WHITE)
text(slide, col2_l, top + Inches(4.5), col2_w, Inches(0.8), "[FEP ddG Bar Chart with Error Bars]", 26, MID_GRAY, False, PP_ALIGN.CENTER)

# FEP interpretation - THE KEY RESULT
rect(slide, col2_l, top + Inches(7.6), col2_w, Inches(3.2), NAVY)
fep_result = """FEBUXOSTAT: ddG = +2.34 +/- 0.26 kcal/mol
Interpretation: 50x weaker binding in R702W | 8 sigma | p < 0.001

BUFADIENOLIDE: ddG = -0.44 +/- 0.37 kcal/mol
Interpretation: No significant change | 1.8 sigma | p = NS"""
text(slide, col2_l + Inches(0.2), top + Inches(7.8), col2_w - Inches(0.4), Inches(2.8), fep_result, 28, WHITE, True, PP_ALIGN.CENTER)

text(slide, col2_l, top + Inches(10.6), col2_w, Inches(0.5), "Fig 2: FEP binding calculations. Positive ddG indicates weaker mutant binding.", 20, DARK_TEXT, False, PP_ALIGN.CENTER)

fep_discuss = """The contrasting results are striking: Febuxostat shows highly significant mutation-sensitive binding while Bufadienolide is unaffected. This ligand-dependent effect suggests the R702W mutation selectively disrupts certain binding modes. Febuxostat's carboxylic acid group likely interacts with residues whose positions are altered by the distant R702W substitution through an allosteric network."""

text(slide, col2_l, top + Inches(11), col2_w, Inches(2.5), fep_discuss, 28, DARK_TEXT)

header(slide, col2_l, top + Inches(13.4), col2_w, "Apo MD Conformational Analysis")

md_intro = """To understand how R702W affects protein dynamics independent of ligand binding, I performed apo molecular dynamics comparing wild-type and mutant NOD2."""

text(slide, col2_l, top + Inches(14), col2_w, Inches(1.5), md_intro, 28, DARK_TEXT)

# MD chart placeholder
rect(slide, col2_l, top + Inches(15.4), col2_w, Inches(4), WHITE)
text(slide, col2_l, top + Inches(17), col2_w, Inches(0.8), "[MD RMSD Distribution Chart]", 24, MID_GRAY, False, PP_ALIGN.CENTER)

md_result = """WT Apo: Mean RMSD 4.05A | 12.6% frames exceed 5A threshold
R702W Apo: Mean RMSD 3.27A | 0% frames exceed 5A threshold"""
text(slide, col2_l, top + Inches(19.2), col2_w, Inches(1.2), md_result, 28, ACCENT, True, PP_ALIGN.CENTER)

text(slide, col2_l, top + Inches(20.2), col2_w, Inches(0.5), "Fig 3: Conformational heterogeneity comparison (n=4 WT, n=3 R702W replicates)", 20, DARK_TEXT, False, PP_ALIGN.CENTER)

md_discuss = """The R702W mutant shows reduced conformational sampling compared to wild-type, with no frames exceeding the 5A RMSD threshold versus 12.6% for WT. This 'rigidification' is consistent with the mutation-sensitive binding observed in FEP: a less dynamic protein may be unable to accommodate certain ligands through induced-fit mechanisms."""

text(slide, col2_l, top + Inches(20.6), col2_w, Inches(2.8), md_discuss, 28, DARK_TEXT)

# ============================================================
# COLUMN 3: STRUCTURAL ANALYSIS + CONCLUSIONS
# ============================================================

header(slide, col3_l, top, col3_w, "Structural Analysis")

struct_text = """The R702W mutation resides in the HD2 domain, approximately 80 angstroms from the LRR ligand binding pocket. This physical separation necessitates an allosteric mechanism to explain the observed binding differences."""

text(slide, col3_l, top + Inches(0.6), col3_w, Inches(2), struct_text, 28, DARK_TEXT)

# Structure figure
rect(slide, col3_l, top + Inches(2.5), col3_w, Inches(5), WHITE)
text(slide, col3_l, top + Inches(4.5), col3_w, Inches(0.8), "[NOD2 Full Structure - Domain Colored]", 22, MID_GRAY, False, PP_ALIGN.CENTER)
text(slide, col3_l, top + Inches(7.3), col3_w, Inches(0.5), "Fig 4: NOD2 architecture. HD2 (orange) contains R702W; LRR (green) contains binding site.", 18, DARK_TEXT, False, PP_ALIGN.CENTER)

# Binding site figure
rect(slide, col3_l, top + Inches(8), col3_w, Inches(4.5), WHITE)
text(slide, col3_l, top + Inches(9.8), col3_w, Inches(0.8), "[Binding Site with H-bonds]", 22, MID_GRAY, False, PP_ALIGN.CENTER)
text(slide, col3_l, top + Inches(12.3), col3_w, Inches(0.5), "Fig 5: Febuxostat in LRR pocket. Yellow dashes = H-bonds to GLU1008, ARG1037.", 18, DARK_TEXT, False, PP_ALIGN.CENTER)

header(slide, col3_l, top + Inches(13), col3_w, "Key Binding Residues (PLIP)")

plip_text = """Protein-Ligand Interaction Profiler analysis identified critical contacts:
- ARG1034: H-bond (3.13A)
- THR1036: H-bond (1.91A)
- ARG1037: H-bond (2.93A)
- Hydrophobic: LEU1007, VAL1009

Contact frequencies from MD: GLU1008 (50-71%), ARG1037 (71-86%), ASP1011 (67-93%)"""

text(slide, col3_l, top + Inches(13.6), col3_w, Inches(3.5), plip_text, 26, DARK_TEXT)

header(slide, col3_l, top + Inches(17), col3_w, "Conclusions")

conclusions = """1. Developed NOD2-Scout ML pipeline (AUC 0.85-0.93) and screened 9,566 compounds against the NOD2 LRR domain.

2. FEP calculations reveal ligand-dependent mutation effects: Febuxostat binding is 50x weaker in R702W (ddG = +2.34 kcal/mol, p<0.001), while Bufadienolide shows no significant change.

3. Apo MD demonstrates reduced conformational sampling in R702W (0% vs 12.6% frames >5A), providing mechanistic context for mutation-sensitive binding.

4. The 80A separation between R702W and the binding site implicates an allosteric mechanism transmitting mutation effects across domains."""

text(slide, col3_l, top + Inches(17.6), col3_w, Inches(5.5), conclusions, 26, DARK_TEXT)

header(slide, col3_l, top + Inches(23), col3_w, "Future Directions")

future = """- Experimental validation: SPR and ITC binding assays
- HDX-MS to probe allosteric communication pathway
- Extend to additional Crohn's variants (G908R, L1007fs)
- Cell-based NOD2 activation reporter assays"""

text(slide, col3_l, top + Inches(23.6), col3_w, Inches(2.5), future, 26, DARK_TEXT)

# ============================================================
# BOTTOM: References
# ============================================================
rect(slide, Inches(0), Inches(34.8), Inches(48), Inches(1.2), NAVY)
refs = "References: Hugot et al. Nature 2001 | Jumper et al. Nature 2021 | Wang et al. JACS 2015 | Figures: PyMOL, PLIP (plip-tool.biotec.tu-dresden.de)"
text(slide, Inches(0.5), Inches(35), Inches(47), Inches(0.8), refs, 20, WHITE, False, PP_ALIGN.CENTER)

# Save
out = r"C:\Users\vasud\nod2-screening-data\isef_figures\v5_winner\ISEF_2026_POSTER_V3.pptx"
prs.save(out)
print(f"SAVED: {out}")
print("V3: Dense text + balanced visuals like Rishab Jain")
