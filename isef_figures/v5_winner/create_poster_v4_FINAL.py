"""
ISEF 2026 Poster V4 FINAL - Fixed spacing + actual images inserted
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(48)
prs.slide_height = Inches(36)
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Colors
NAVY = RGBColor(0x1a, 0x36, 0x5d)
DARK_TEXT = RGBColor(0x1f, 0x2d, 0x3d)
ACCENT = RGBColor(0xc0, 0x39, 0x2b)
LIGHT_BG = RGBColor(0xf5, 0xf6, 0xf8)
WHITE = RGBColor(255, 255, 255)
MID_GRAY = RGBColor(0x64, 0x74, 0x8b)

img_folder = r"C:\Users\vasud\nod2-screening-data\isef_figures\v5_winner\ppt_photos"

def rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()

def text(slide, l, t, w, h, txt, size, color=DARK_TEXT, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = txt
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align

def header(slide, l, t, w, txt):
    rect(slide, l, t, w, Inches(0.5), NAVY)
    text(slide, l + Inches(0.1), t + Inches(0.05), w - Inches(0.2), Inches(0.45), txt.upper(), 28, WHITE, True)

def add_image(slide, path, l, t, w, h):
    if os.path.exists(path):
        slide.shapes.add_picture(path, l, t, w, h)
    else:
        rect(slide, l, t, w, h, WHITE)
        text(slide, l, t + h/2 - Inches(0.3), w, Inches(0.6), f"[Missing: {os.path.basename(path)}]", 18, MID_GRAY, False, PP_ALIGN.CENTER)

# Background
rect(slide, Inches(0), Inches(0), Inches(48), Inches(36), LIGHT_BG)

# ============================================================
# TITLE BANNER - 2.8" tall
# ============================================================
rect(slide, Inches(0), Inches(0), Inches(48), Inches(2.8), NAVY)

text(slide, Inches(0.5), Inches(0.3), Inches(47), Inches(1.2),
     "Deep Learning Guided Discovery and FEP Validation of NOD2 Binders",
     64, WHITE, True, PP_ALIGN.CENTER)

text(slide, Inches(0.5), Inches(1.3), Inches(47), Inches(0.6),
     "for the R702W Crohn's Disease Variant",
     48, WHITE, False, PP_ALIGN.CENTER)

text(slide, Inches(0.5), Inches(2), Inches(47), Inches(0.5),
     "Tanmay Vasudeva  |  Texas Virtual Academy at Hallsville  |  ISEF 2026",
     28, WHITE, False, PP_ALIGN.CENTER)

# Key finding - separate box below title
rect(slide, Inches(6), Inches(3), Inches(36), Inches(0.7), ACCENT)
text(slide, Inches(6.2), Inches(3.1), Inches(35.6), Inches(0.5),
     "KEY FINDING: Febuxostat shows 50x weaker binding in R702W mutant (ddG = +2.34 kcal/mol, p<0.001)",
     26, WHITE, True, PP_ALIGN.CENTER)

# ============================================================
# 3-COLUMN LAYOUT - More space
# ============================================================
margin = Inches(0.5)
gap = Inches(0.4)
top = Inches(4)

col1_l = margin
col1_w = Inches(14)

col2_l = col1_l + col1_w + gap
col2_w = Inches(16.2)

col3_l = col2_l + col2_w + gap
col3_w = Inches(16.2)

# ============================================================
# COLUMN 1: INTRO + METHODS
# ============================================================

y = top
header(slide, col1_l, y, col1_w, "Research Question")
y += Inches(0.55)

q = """What mechanisms explain how the NOD2 R702W variant alters drug binding? I developed a pipeline combining ML screening with FEP calculations to quantify mutation-sensitive binding for FDA compounds."""
text(slide, col1_l, y, col1_w, Inches(1.8), q, 26, DARK_TEXT)
y += Inches(1.9)

header(slide, col1_l, y, col1_w, "Background")
y += Inches(0.55)

bg = """NOD2 detects bacterial peptidoglycans and initiates NF-kB signaling. R702W causes 15-25% of Crohn's disease cases. The mutation is in HD2 domain, 80A from the LRR binding site. No targeted therapies exist for R702W carriers. Drug repurposing offers faster path to therapy."""
text(slide, col1_l, y, col1_w, Inches(2.5), bg, 26, DARK_TEXT)
y += Inches(2.6)

header(slide, col1_l, y, col1_w, "Methods")
y += Inches(0.55)

methods = """1. Structure: AlphaFold2 NOD2 (Q9HC29), LRR domain extracted (res 744-1040)

2. Screening: 9,566 FDA + natural compounds via GNINA CNN docking

3. ML Scoring: NOD2-Scout classifier, AUC 0.85-0.93 (scaffold-split CV)

4. MD Simulation: OpenMM, AMBER14, 520 ns total, 310K, TIP3P-FB

5. FEP Binding: 120 windows per ligand, thermodynamic cycle, ddG = dG(mut) - dG(WT)"""
text(slide, col1_l, y, col1_w, Inches(4.5), methods, 24, DARK_TEXT)
y += Inches(4.6)

# Pipeline image
add_image(slide, os.path.join(img_folder, "pipeline.png"), col1_l, y, col1_w, Inches(4))
y += Inches(4.1)
text(slide, col1_l, y, col1_w, Inches(0.4), "Fig 1: Computational pipeline", 18, DARK_TEXT, False, PP_ALIGN.CENTER)
y += Inches(0.5)

header(slide, col1_l, y, col1_w, "Limitations")
y += Inches(0.55)

limits = """FEP on 2 ligands (compute cost). Predictions need experimental validation. Apo MD only."""
text(slide, col1_l, y, col1_w, Inches(1.2), limits, 24, DARK_TEXT)

# ============================================================
# COLUMN 2: RESULTS
# ============================================================

y = top
header(slide, col2_l, y, col2_w, "FEP Binding Free Energy")
y += Inches(0.55)

fep_intro = """FEP is the gold standard for binding affinity. Positive ddG = weaker binding in mutant."""
text(slide, col2_l, y, col2_w, Inches(0.9), fep_intro, 26, DARK_TEXT)
y += Inches(1)

# FEP chart
add_image(slide, os.path.join(img_folder, "fep_chart.png"), col2_l, y, col2_w, Inches(4.5))
y += Inches(4.6)
text(slide, col2_l, y, col2_w, Inches(0.4), "Fig 2: FEP binding calculations", 18, DARK_TEXT, False, PP_ALIGN.CENTER)
y += Inches(0.5)

# FEP results box
rect(slide, col2_l, y, col2_w, Inches(2), NAVY)
fep_data = """FEBUXOSTAT: ddG = +2.34 kcal/mol | 50x weaker | p<0.001 (8s)
BUFADIENOLIDE: ddG = -0.44 kcal/mol | No change | NS (1.8s)"""
text(slide, col2_l + Inches(0.2), y + Inches(0.3), col2_w - Inches(0.4), Inches(1.5), fep_data, 26, WHITE, True, PP_ALIGN.CENTER)
y += Inches(2.2)

fep_discuss = """Febuxostat shows significant mutation-sensitive binding while Bufadienolide is unaffected. This ligand-dependent effect suggests R702W selectively disrupts certain binding modes through allosteric communication."""
text(slide, col2_l, y, col2_w, Inches(1.8), fep_discuss, 26, DARK_TEXT)
y += Inches(2)

header(slide, col2_l, y, col2_w, "Apo MD Conformational Analysis")
y += Inches(0.55)

md_intro = """MD without ligand reveals how R702W affects protein dynamics."""
text(slide, col2_l, y, col2_w, Inches(0.7), md_intro, 26, DARK_TEXT)
y += Inches(0.8)

# MD chart
add_image(slide, os.path.join(img_folder, "md_rmsd.png"), col2_l, y, Inches(7.8), Inches(4))
add_image(slide, os.path.join(img_folder, "heterogeneity.png"), col2_l + Inches(8.2), y, Inches(7.8), Inches(4))
y += Inches(4.1)
text(slide, col2_l, y, col2_w, Inches(0.4), "Fig 3: RMSD traces and conformational heterogeneity", 18, DARK_TEXT, False, PP_ALIGN.CENTER)
y += Inches(0.5)

md_result = """WT: 12.6% frames >5A (higher flexibility) | R702W: 0% frames >5A (reduced sampling)"""
text(slide, col2_l, y, col2_w, Inches(0.7), md_result, 26, ACCENT, True, PP_ALIGN.CENTER)
y += Inches(0.8)

md_discuss = """R702W shows reduced conformational sampling - the mutant is 'rigidified'. This may prevent induced-fit binding, explaining mutation-sensitive effects."""
text(slide, col2_l, y, col2_w, Inches(1.5), md_discuss, 26, DARK_TEXT)

# ============================================================
# COLUMN 3: STRUCTURE + CONCLUSIONS
# ============================================================

y = top
header(slide, col3_l, y, col3_w, "Structural Analysis")
y += Inches(0.55)

struct = """R702W is 80A from binding site - requires allosteric mechanism."""
text(slide, col3_l, y, col3_w, Inches(0.7), struct, 26, DARK_TEXT)
y += Inches(0.8)

# Full structure
add_image(slide, os.path.join(img_folder, "figA1_full_domains_white.png"), col3_l, y, col3_w, Inches(5))
y += Inches(5.1)
text(slide, col3_l, y, col3_w, Inches(0.5), "Fig 4: NOD2 domains. HD2 (orange)=R702W site, LRR (green)=binding", 18, DARK_TEXT, False, PP_ALIGN.CENTER)
y += Inches(0.6)

# Binding site
add_image(slide, os.path.join(img_folder, "figB1_binding_hbonds_white.png"), col3_l, y, col3_w, Inches(4.5))
y += Inches(4.6)
text(slide, col3_l, y, col3_w, Inches(0.5), "Fig 5: Febuxostat binding. H-bonds to GLU1008, ARG1037", 18, DARK_TEXT, False, PP_ALIGN.CENTER)
y += Inches(0.6)

header(slide, col3_l, y, col3_w, "Conclusions")
y += Inches(0.55)

conclusions = """1. Built NOD2-Scout ML pipeline (AUC 0.85-0.93), screened 9,566 compounds

2. FEP shows ligand-dependent mutation effects:
   - Febuxostat: 50x weaker in R702W (p<0.001)
   - Bufadienolide: No significant change

3. Apo MD: R702W has reduced conformational sampling (0% vs 12.6% >5A)

4. 80A distance implies allosteric mechanism"""
text(slide, col3_l, y, col3_w, Inches(4), conclusions, 24, DARK_TEXT)
y += Inches(4.2)

header(slide, col3_l, y, col3_w, "Future Work")
y += Inches(0.55)

future = """- SPR/ITC binding validation
- HDX-MS for allosteric pathway
- Test G908R, L1007fs mutations
- Cell-based NOD2 reporter assays"""
text(slide, col3_l, y, col3_w, Inches(2), future, 24, DARK_TEXT)

# ============================================================
# BOTTOM: References
# ============================================================
rect(slide, Inches(0), Inches(35), Inches(48), Inches(1), NAVY)
refs = "References: Hugot et al. Nature 2001 | Jumper et al. Nature 2021 | Wang et al. JACS 2015 | Figures by author (PyMOL, matplotlib)"
text(slide, Inches(0.5), Inches(35.15), Inches(47), Inches(0.7), refs, 20, WHITE, False, PP_ALIGN.CENTER)

# Save
out = r"C:\Users\vasud\nod2-screening-data\isef_figures\v5_winner\ISEF_2026_POSTER_V4.pptx"
prs.save(out)
print(f"SAVED: {out}")
print("V4: Fixed spacing + actual images inserted")
