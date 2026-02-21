"""
ISEF 2026 Poster - V2 FIXED
- NO Abstract (ISEF rule)
- Proper color scheme (not tacky)
- Bigger text
- Real sections with content
- Space for DATA CHARTS + PyMOL figures
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Create 48x36 landscape poster
prs = Presentation()
prs.slide_width = Inches(48)
prs.slide_height = Inches(36)
slide = prs.slides.add_slide(prs.slide_layouts[6])

# ============================================================
# COLOR SCHEME - Option B (Clean Biomedical) - NOT TACKY
# ============================================================
DEEP_BLUE = RGBColor(0x14, 0x3A, 0x5A)      # Headers
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)      # Body text
MUTED_BLUE = RGBColor(0x4C, 0x6A, 0x87)     # Subheads
ACCENT_ORANGE = RGBColor(0xE0, 0x7A, 0x5F)  # Key results highlight
OFF_WHITE = RGBColor(0xFA, 0xFA, 0xFB)      # Background
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(0xE8, 0xEB, 0xEE)     # Panel backgrounds

def add_rect(slide, left, top, width, height, color, outline=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if not outline:
        shape.line.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, size, color=DARK_GRAY, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return box

def section_header(slide, left, top, width, text):
    """Deep blue header bar"""
    add_rect(slide, left, top, width, Inches(0.6), DEEP_BLUE)
    add_text(slide, left + Inches(0.15), top + Inches(0.08), width, Inches(0.5),
             text.upper(), 36, WHITE, bold=True)

# ============================================================
# BACKGROUND - Off-white (professional, not dark)
# ============================================================
add_rect(slide, Inches(0), Inches(0), Inches(48), Inches(36), OFF_WHITE)

# ============================================================
# TITLE BANNER - Deep Blue, spans full width
# ============================================================
add_rect(slide, Inches(0), Inches(0), Inches(48), Inches(4), DEEP_BLUE)

# Title - BIG (140pt)
add_text(slide, Inches(1), Inches(0.5), Inches(46), Inches(1.5),
         "Deep Learning Guided Discovery and FEP Validation of NOD2 Binders for the R702W Crohn's Variant",
         80, WHITE, bold=True, align=PP_ALIGN.CENTER)

# Subtitle / Key Finding - ONE LINE
add_text(slide, Inches(1), Inches(2), Inches(46), Inches(0.8),
         "Key Finding: Febuxostat shows 50× weaker binding in R702W mutant (ΔΔG = +2.34 kcal/mol, p<0.001)",
         36, ACCENT_ORANGE, bold=True, align=PP_ALIGN.CENTER)

# Author
add_text(slide, Inches(1), Inches(3), Inches(46), Inches(0.7),
         "Tanmay Vasudeva  |  Texas Virtual Academy at Hallsville  |  ISEF 2026",
         32, WHITE, align=PP_ALIGN.CENTER)

# ============================================================
# 3-COLUMN LAYOUT
# Column 1: 12" (Intro + Methods)
# Column 2: 17" (Structural Figures)
# Column 3: 17" (Data Results)
# ============================================================
margin = Inches(0.75)
col_gap = Inches(0.5)
top_start = Inches(4.5)

col1_left = margin
col1_width = Inches(12)

col2_left = col1_left + col1_width + col_gap
col2_width = Inches(16.5)

col3_left = col2_left + col2_width + col_gap
col3_width = Inches(16.5)

# ============================================================
# COLUMN 1: INTRODUCTION + METHODS
# ============================================================

# RESEARCH QUESTION
section_header(slide, col1_left, top_start, col1_width, "Research Question")
question_text = """Can computational drug repurposing identify FDA-approved compounds that bind the NOD2 LRR domain, and how does the pathogenic R702W mutation affect their binding affinity?

The R702W mutation causes 15-25% of Crohn's disease cases, yet no targeted therapies exist. This project combines machine learning screening with rigorous free energy perturbation (FEP) calculations to quantify mutation-sensitive binding."""

add_text(slide, col1_left + Inches(0.1), top_start + Inches(0.7), col1_width - Inches(0.2), Inches(4),
         question_text, 26, DARK_GRAY)

# BACKGROUND
section_header(slide, col1_left, top_start + Inches(4.5), col1_width, "Background")
background_text = """• NOD2 is an innate immune receptor that detects bacterial peptidoglycans
• R702W mutation disrupts LRR domain function, causing impaired bacterial sensing
• Located in HD2 domain, ~80Å from LRR binding site (allosteric effect)
• Current treatments (anti-TNF, steroids) don't address root cause
• Drug repurposing offers faster path to therapy than de novo design"""

add_text(slide, col1_left + Inches(0.1), top_start + Inches(5.2), col1_width - Inches(0.2), Inches(4.5),
         background_text, 26, DARK_GRAY)

# METHODS
section_header(slide, col1_left, top_start + Inches(9.5), col1_width, "Computational Methods")
methods_text = """1. TARGET PREPARATION
   • AlphaFold2 structure of NOD2 (Q9HC29)
   • LRR domain extracted (residues 744-1040)

2. VIRTUAL SCREENING
   • 9,566 FDA + natural product compounds
   • GNINA CNN-based docking
   • NOD2-Scout ML model (AUC: 0.85-0.93)

3. MOLECULAR DYNAMICS
   • OpenMM engine, AMBER14 + TIP3P-FB
   • 520 ns total simulation time
   • Apo WT vs R702W comparison

4. FEP BINDING CALCULATIONS
   • Gold-standard thermodynamic method
   • 120 windows (80 complex + 40 solvent)
   • ΔΔG = ΔG(R702W) - ΔG(WT)"""

add_text(slide, col1_left + Inches(0.1), top_start + Inches(10.2), col1_width - Inches(0.2), Inches(9),
         methods_text, 24, DARK_GRAY)

# Pipeline placeholder
add_rect(slide, col1_left, top_start + Inches(19.5), col1_width, Inches(5), LIGHT_GRAY)
add_text(slide, col1_left, top_start + Inches(21.5), col1_width, Inches(1),
         "[INSERT: Pipeline Flowchart]", 24, MUTED_BLUE, align=PP_ALIGN.CENTER)
add_text(slide, col1_left, top_start + Inches(24), col1_width, Inches(0.6),
         "Figure 1: Computational pipeline from screening to FEP validation", 20, DARK_GRAY, align=PP_ALIGN.CENTER)

# LIMITATIONS
section_header(slide, col1_left, top_start + Inches(25), col1_width, "Limitations")
limits_text = """• FEP performed on 2 representative ligands (compute cost)
• Computational predictions require experimental validation
• Apo MD only; ligand-bound dynamics not compared"""

add_text(slide, col1_left + Inches(0.1), top_start + Inches(25.7), col1_width - Inches(0.2), Inches(2.5),
         limits_text, 24, DARK_GRAY)

# ============================================================
# COLUMN 2: STRUCTURAL FIGURES (PyMOL + PLIP)
# ============================================================

section_header(slide, col2_left, top_start, col2_width, "Structural Analysis")

# Full NOD2 structure - BIG
add_rect(slide, col2_left, top_start + Inches(0.7), col2_width, Inches(8), LIGHT_GRAY)
add_text(slide, col2_left, top_start + Inches(4), col2_width, Inches(1),
         "[INSERT: figA1_full_domains_white.png]", 28, MUTED_BLUE, align=PP_ALIGN.CENTER)
add_text(slide, col2_left, top_start + Inches(8.3), col2_width, Inches(1),
         "Figure 2: NOD2 domain architecture. CARD (gray), NBD (blue), HD2 (orange, contains R702W), LRR (green, binding site). Red sphere = R702W mutation, magenta = binding pocket.",
         22, DARK_GRAY, align=PP_ALIGN.CENTER)

# R702W closeup + Binding site side by side
section_header(slide, col2_left, top_start + Inches(9.8), col2_width, "Mutation Site & Binding Pocket")

half_width = (col2_width - Inches(0.3)) / 2

# R702W closeup
add_rect(slide, col2_left, top_start + Inches(10.5), half_width, Inches(6.5), LIGHT_GRAY)
add_text(slide, col2_left, top_start + Inches(13), half_width, Inches(1),
         "[INSERT: figA3_r702w_closeup.png]", 24, MUTED_BLUE, align=PP_ALIGN.CENTER)
add_text(slide, col2_left, top_start + Inches(16.7), half_width, Inches(0.8),
         "Figure 3: R702W mutation site in HD2 domain", 20, DARK_GRAY, align=PP_ALIGN.CENTER)

# Binding site
add_rect(slide, col2_left + half_width + Inches(0.3), top_start + Inches(10.5), half_width, Inches(6.5), LIGHT_GRAY)
add_text(slide, col2_left + half_width + Inches(0.3), top_start + Inches(13), half_width, Inches(1),
         "[INSERT: figB1_binding_hbonds_white.png]", 24, MUTED_BLUE, align=PP_ALIGN.CENTER)
add_text(slide, col2_left + half_width + Inches(0.3), top_start + Inches(16.7), half_width, Inches(0.8),
         "Figure 4: Febuxostat binding. H-bonds (yellow) shown.", 20, DARK_GRAY, align=PP_ALIGN.CENTER)

# PLIP 2D diagram
section_header(slide, col2_left, top_start + Inches(18), col2_width, "Protein-Ligand Interactions (PLIP)")

add_rect(slide, col2_left, top_start + Inches(18.7), col2_width, Inches(7), LIGHT_GRAY)
add_text(slide, col2_left, top_start + Inches(21.5), col2_width, Inches(1),
         "[INSERT: PLIP 2D Interaction Diagram]", 28, MUTED_BLUE, align=PP_ALIGN.CENTER)

plip_caption = """Figure 5: PLIP analysis confirms key contacts. H-bonds: ARG1034 (3.13Å), THR1036 (1.91Å), ARG1037 (2.93Å).
Hydrophobic contacts with LEU1007, VAL1009. Binding pocket is 80Å from R702W (allosteric mechanism)."""
add_text(slide, col2_left, top_start + Inches(25.4), col2_width, Inches(1.2),
         plip_caption, 22, DARK_GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# COLUMN 3: DATA RESULTS + CONCLUSIONS
# ============================================================

section_header(slide, col3_left, top_start, col3_width, "FEP Binding Free Energy Results")

# FEP bar chart - THE MAIN RESULT
add_rect(slide, col3_left, top_start + Inches(0.7), col3_width, Inches(7), LIGHT_GRAY)
add_text(slide, col3_left, top_start + Inches(3.5), col3_width, Inches(1),
         "[INSERT: FEP ΔΔG Bar Chart]", 28, MUTED_BLUE, align=PP_ALIGN.CENTER)

# FEP data callouts - BIG AND BOLD
fep_box = add_rect(slide, col3_left + Inches(0.5), top_start + Inches(7.5), col3_width - Inches(1), Inches(3.5), DEEP_BLUE)

fep_data = """FEBUXOSTAT: ΔΔG = +2.34 ± 0.26 kcal/mol
→ 50× WEAKER in R702W mutant | p < 0.001 (8σ) ★★★

BUFADIENOLIDE: ΔΔG = -0.44 ± 0.37 kcal/mol
→ NO SIGNIFICANT CHANGE | p = NS (1.8σ)"""

add_text(slide, col3_left + Inches(0.7), top_start + Inches(7.8), col3_width - Inches(1.4), Inches(3),
         fep_data, 28, WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, col3_left, top_start + Inches(10.8), col3_width, Inches(0.6),
         "Figure 6: FEP binding calculations. Positive ΔΔG = weaker binding in mutant.", 20, DARK_GRAY, align=PP_ALIGN.CENTER)

# MD Results
section_header(slide, col3_left, top_start + Inches(11.5), col3_width, "Apo MD Conformational Analysis")

add_rect(slide, col3_left, top_start + Inches(12.2), col3_width, Inches(5), LIGHT_GRAY)
add_text(slide, col3_left, top_start + Inches(14), col3_width, Inches(1),
         "[INSERT: MD RMSD / Conformational Chart]", 24, MUTED_BLUE, align=PP_ALIGN.CENTER)

md_data = """WT Apo: 12.6% frames with RMSD >5Å (higher conformational heterogeneity)
R702W Apo: 0% frames with RMSD >5Å (reduced conformational sampling)"""

add_text(slide, col3_left + Inches(0.2), top_start + Inches(17), col3_width - Inches(0.4), Inches(1.5),
         md_data, 26, DARK_GRAY, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, col3_left, top_start + Inches(18.2), col3_width, Inches(0.6),
         "Figure 7: Apo MD shows reduced conformational sampling in R702W mutant", 20, DARK_GRAY, align=PP_ALIGN.CENTER)

# CONCLUSIONS
section_header(slide, col3_left, top_start + Inches(19.2), col3_width, "Conclusions")

conclusions = """1. Built ML screening pipeline (NOD2-Scout, AUC 0.85-0.93) and screened 9,566 compounds against NOD2 LRR domain

2. FEP calculations quantify mutation-sensitive binding:
   • Febuxostat: +2.34 kcal/mol weaker in R702W (50×, p<0.001)
   • Bufadienolide: No significant change between WT and R702W

3. Apo MD reveals reduced conformational sampling in R702W (0% vs 12.6% frames >5Å), consistent with mutation-sensitive binding

4. R702W is 80Å from binding site, suggesting allosteric mechanism for mutation-dependent binding differences"""

add_text(slide, col3_left + Inches(0.1), top_start + Inches(19.9), col3_width - Inches(0.2), Inches(5.5),
         conclusions, 24, DARK_GRAY)

# FUTURE WORK
section_header(slide, col3_left, top_start + Inches(25.5), col3_width, "Future Directions")

future = """• Experimental validation: SPR, ITC binding assays
• HDX-MS to probe conformational dynamics
• Test additional Crohn's mutations (G908R, L1007fs)
• Cell-based NOD2 activation reporter assays"""

add_text(slide, col3_left + Inches(0.1), top_start + Inches(26.2), col3_width - Inches(0.2), Inches(2.5),
         future, 24, DARK_GRAY)

# ============================================================
# BOTTOM: REFERENCES
# ============================================================
add_rect(slide, Inches(0), Inches(34.5), Inches(48), Inches(1.5), DEEP_BLUE)
refs = "References: Hugot et al. Nature 2001 | Jumper et al. Nature 2021 (AlphaFold) | Wang et al. JACS 2015 (FEP) | All figures by author using PyMOL, PLIP, matplotlib"
add_text(slide, Inches(1), Inches(34.7), Inches(46), Inches(1), refs, 20, WHITE, align=PP_ALIGN.CENTER)

# Save
output = r"C:\Users\vasud\nod2-screening-data\isef_figures\v5_winner\ISEF_2026_POSTER_V2.pptx"
prs.save(output)
print(f"SAVED: {output}")
print("\n✓ NO Abstract section")
print("✓ Professional color scheme (Deep Blue + Orange accent)")
print("✓ Bigger text (24-36pt body)")
print("✓ Real content with data")
print("✓ Space for 7 figures + charts")
print("\nNow insert your actual figures and data charts!")
