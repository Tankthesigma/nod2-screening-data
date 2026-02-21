"""
ISEF 2026 Poster - Rishab Jain Style
Dense, Dark Blue, Visually Aggressive
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Create presentation - 48" x 36" landscape
prs = Presentation()
prs.slide_width = Inches(48)
prs.slide_height = Inches(36)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

# Colors - Rishab Jain Style
DARK_BLUE = RGBColor(0x0A, 0x1F, 0x44)  # Deep navy
HEADER_BLUE = RGBColor(0x1E, 0x3A, 0x5F)  # Section headers
ACCENT_GREEN = RGBColor(0x2E, 0x7D, 0x32)  # For highlights
ACCENT_ORANGE = RGBColor(0xF5, 0x7C, 0x00)  # For key results
WHITE = RGBColor(255, 255, 255)
LIGHT_BLUE = RGBColor(0x42, 0x6A, 0x8C)

def add_rectangle(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return box

def add_section_header(slide, left, top, width, text):
    """Green header bar like Rishab's"""
    # Green background bar
    bar = add_rectangle(slide, left, top, width, Inches(0.5), ACCENT_GREEN)
    # White text
    add_text_box(slide, left + Inches(0.1), top + Inches(0.05), width - Inches(0.2), Inches(0.4),
                 text.upper(), 24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ============================================================
# BACKGROUND - Full dark blue
# ============================================================
add_rectangle(slide, Inches(0), Inches(0), Inches(48), Inches(36), DARK_BLUE)

# ============================================================
# HEADER BANNER - Title spanning full width
# ============================================================
add_rectangle(slide, Inches(0), Inches(0), Inches(48), Inches(3.5), HEADER_BLUE)

# Title
add_text_box(slide, Inches(1), Inches(0.4), Inches(46), Inches(1.2),
             "DEEP LEARNING GUIDED DISCOVERY AND FEP VALIDATION OF NOD2 BINDERS",
             72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Subtitle
add_text_box(slide, Inches(1), Inches(1.5), Inches(46), Inches(0.8),
             "Targeting the R702W Crohn's Disease Variant with Structure-Based Drug Repurposing",
             36, bold=False, color=WHITE, align=PP_ALIGN.CENTER)

# Author
add_text_box(slide, Inches(1), Inches(2.4), Inches(46), Inches(0.6),
             "Tanmay Vasudeva  •  Texas Virtual Academy at Hallsville",
             28, bold=False, color=WHITE, align=PP_ALIGN.CENTER)

# ============================================================
# COLUMN LAYOUT (3 columns: 25% / 40% / 35%)
# ============================================================
col1_left = Inches(0.5)
col1_width = Inches(11)  # ~25%

col2_left = Inches(12)
col2_width = Inches(18)  # ~40%

col3_left = Inches(31)
col3_width = Inches(16)  # ~35%

content_top = Inches(4)

# ============================================================
# LEFT COLUMN - The Setup
# ============================================================

# VISUAL ABSTRACT (top left box)
add_section_header(slide, col1_left, content_top, col1_width, "Visual Abstract")

# Visual abstract box (white background for contrast)
va_box = add_rectangle(slide, col1_left, content_top + Inches(0.6), col1_width, Inches(4), WHITE)

# Text inside visual abstract
add_text_box(slide, col1_left + Inches(0.3), content_top + Inches(0.8), col1_width - Inches(0.6), Inches(3.5),
             "NOD2 Mutation (R702W)\n     ↓\nReduced LRR Stability\n     ↓\nML Screening (9,566 compounds)\n     ↓\nFEP Validation\n     ↓\nFebuxostat: 50× weaker in mutant\nBuffadienolide: No change",
             24, bold=False, color=DARK_BLUE, align=PP_ALIGN.CENTER)

# INTRODUCTION
add_section_header(slide, col1_left, content_top + Inches(5), col1_width, "Introduction")

intro_text = """• Crohn's disease affects 3M+ people globally
• NOD2 R702W mutation causes 15-25% of cases
• R702W disrupts LRR domain sensing function
• No targeted therapies exist for this variant
• Hypothesis: Drug repurposing can restore function"""

add_text_box(slide, col1_left + Inches(0.2), content_top + Inches(5.6), col1_width - Inches(0.4), Inches(4),
             intro_text, 22, color=WHITE)

# PROBLEM FRAMING
add_section_header(slide, col1_left, content_top + Inches(10), col1_width, "Problem Framing")

problem_text = """• R702W is 80Å from binding site (allosteric)
• FEP shows mutation-sensitive binding
• Need: Identify compounds that bind despite mutation"""

add_text_box(slide, col1_left + Inches(0.2), content_top + Inches(10.6), col1_width - Inches(0.4), Inches(3),
             problem_text, 22, color=WHITE)

# Image placeholder box for figA1
add_rectangle(slide, col1_left, content_top + Inches(13.5), col1_width, Inches(6), LIGHT_BLUE)
add_text_box(slide, col1_left, content_top + Inches(15.5), col1_width, Inches(1),
             "[INSERT: figA1_full_domains_white.png]", 20, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, col1_left, content_top + Inches(19.2), col1_width, Inches(0.8),
             "Fig 1: NOD2 Domain Architecture. R702W (red) in HD2 domain.", 16, color=WHITE, align=PP_ALIGN.CENTER)

# RESEARCH OBJECTIVES
add_section_header(slide, col1_left, content_top + Inches(20.5), col1_width, "Research Objectives")

objectives_text = """1. Screen FDA library against NOD2 LRR
2. Validate hits with FEP binding calculations
3. Compare WT vs R702W binding affinity
4. Identify mutation-sensitive compounds"""

add_text_box(slide, col1_left + Inches(0.2), content_top + Inches(21.1), col1_width - Inches(0.4), Inches(3.5),
             objectives_text, 22, color=WHITE)

# ============================================================
# MIDDLE COLUMN - The Work (Engineering Methodology)
# ============================================================

add_section_header(slide, col2_left, content_top, col2_width, "Computational Pipeline & Methodology")

# Pipeline steps boxes
pipeline_y = content_top + Inches(0.7)
step_width = Inches(4)
step_height = Inches(2.5)

steps = [
    ("Step 1\nTarget Prep", "AlphaFold2\nNOD2 LRR\nStructure"),
    ("Step 2\nML Screening", "NOD2-Scout\nAUC: 0.85-0.93\n9,566 compounds"),
    ("Step 3\nMD Simulation", "520 ns total\nOpenMM\nAMBER14"),
    ("Step 4\nFEP Binding", "Gold Standard\n120 windows\nΔΔG calculation")
]

for i, (title, content) in enumerate(steps):
    x = col2_left + Inches(0.5) + i * (step_width + Inches(0.3))
    # Box
    add_rectangle(slide, x, pipeline_y, step_width, step_height, LIGHT_BLUE)
    # Title
    add_text_box(slide, x, pipeline_y + Inches(0.1), step_width, Inches(0.8), title, 20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Content
    add_text_box(slide, x, pipeline_y + Inches(1), step_width, Inches(1.4), content, 18, color=WHITE, align=PP_ALIGN.CENTER)

# RESULTS & IN-SILICO VALIDATION (Big section)
add_section_header(slide, col2_left, content_top + Inches(3.5), col2_width, "Results & In-Silico Validation")

# FEP Results - THE CENTERPIECE
add_rectangle(slide, col2_left, content_top + Inches(4.2), col2_width, Inches(8), LIGHT_BLUE)
add_text_box(slide, col2_left, content_top + Inches(4.5), col2_width, Inches(1),
             "[INSERT: FEP BAR CHART - BIGGEST FIGURE]", 24, color=WHITE, align=PP_ALIGN.CENTER)

# FEP Key callouts
fep_results = """
┌─────────────────────────────────────────────────────────────┐
│  FEBUXOSTAT                    │  BUFADIENOLIDE            │
│  ΔΔG = +2.34 kcal/mol          │  ΔΔG = -0.44 kcal/mol     │
│  50× WEAKER in R702W           │  NO SIGNIFICANT CHANGE    │
│  p < 0.001 (8σ) ***            │  NS (1.8σ)                │
└─────────────────────────────────────────────────────────────┘
"""
add_text_box(slide, col2_left + Inches(1), content_top + Inches(9), col2_width - Inches(2), Inches(3),
             fep_results, 22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text_box(slide, col2_left, content_top + Inches(11.8), col2_width, Inches(0.6),
             "Fig 3: FEP Binding Free Energy. Mutation-sensitive binding for Febuxostat.", 16, color=WHITE, align=PP_ALIGN.CENTER)

# Binding Site Images
add_section_header(slide, col2_left, content_top + Inches(12.8), col2_width, "Molecular Binding Analysis")

# Two image placeholders side by side
img_width = Inches(8.5)
add_rectangle(slide, col2_left + Inches(0.25), content_top + Inches(13.5), img_width, Inches(6.5), LIGHT_BLUE)
add_text_box(slide, col2_left + Inches(0.25), content_top + Inches(16), img_width, Inches(1),
             "[INSERT: figB1_binding_hbonds_white.png]", 18, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, col2_left + Inches(0.25), content_top + Inches(19.7), img_width, Inches(0.6),
             "Fig 4: Febuxostat Binding. H-bonds (yellow) to key residues.", 16, color=WHITE, align=PP_ALIGN.CENTER)

add_rectangle(slide, col2_left + Inches(9), content_top + Inches(13.5), img_width, Inches(6.5), LIGHT_BLUE)
add_text_box(slide, col2_left + Inches(9), content_top + Inches(16), img_width, Inches(1),
             "[INSERT: PLIP 2D Interaction Diagram]", 18, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, col2_left + Inches(9), content_top + Inches(19.7), img_width, Inches(0.6),
             "Fig 5: PLIP Analysis. ARG1034, THR1036, ARG1037 contacts.", 16, color=WHITE, align=PP_ALIGN.CENTER)

# MD Results
add_section_header(slide, col2_left, content_top + Inches(20.8), col2_width, "Apo MD Dynamics")

md_text = """Conformational Sampling Analysis (Apo MD, no ligand):
• WT: 12.6% frames >5Å RMSD (higher heterogeneity)
• R702W: 0% frames >5Å RMSD (reduced sampling)
• Consistent with mutation-sensitive binding observed in FEP"""

add_text_box(slide, col2_left + Inches(0.3), content_top + Inches(21.4), col2_width - Inches(0.6), Inches(3),
             md_text, 22, color=WHITE)

# ============================================================
# RIGHT COLUMN - The Proof
# ============================================================

add_section_header(slide, col3_left, content_top, col3_width, "Statistical Analysis")

stats_text = """FEP Statistical Validation:
• Febuxostat: 8σ significance (p < 0.001)
• Error propagation: SEM from 3 replicates
• 120 FEP windows (80 complex + 40 solvent)

Key Binding Residues (PLIP confirmed):
• GLU1008: 50-71% contact frequency
• ARG1037: 71-86% contact frequency
• ASP1011: 67-93% contact frequency
• ARG1034: 82-88% contact frequency"""

add_text_box(slide, col3_left + Inches(0.2), content_top + Inches(0.6), col3_width - Inches(0.4), Inches(5),
             stats_text, 20, color=WHITE)

# R702W Closeup image
add_rectangle(slide, col3_left, content_top + Inches(6), col3_width, Inches(5.5), LIGHT_BLUE)
add_text_box(slide, col3_left, content_top + Inches(8), col3_width, Inches(1),
             "[INSERT: figA3_r702w_closeup.png]", 18, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, col3_left, content_top + Inches(11.2), col3_width, Inches(0.6),
             "Fig 6: R702W Mutation Site. 80Å from binding pocket (allosteric).", 16, color=WHITE, align=PP_ALIGN.CENTER)

# CONCLUSIONS
add_section_header(slide, col3_left, content_top + Inches(12), col3_width, "Conclusions")

conclusions_text = """1. Built ML pipeline screening 9,566 compounds
   against NOD2 LRR (AUC 0.85-0.93)

2. FEP quantifies mutation-sensitive binding:
   • Febuxostat: +2.34 kcal/mol (50× weaker, p<0.001)
   • Bufadienolide: No significant change

3. Apo MD shows reduced conformational sampling
   in R702W (0% vs 12.6% frames >5Å)

4. Ligand-dependent binding differences suggest
   allosteric mechanism from distant mutation"""

add_text_box(slide, col3_left + Inches(0.2), content_top + Inches(12.6), col3_width - Inches(0.4), Inches(6),
             conclusions_text, 20, color=WHITE)

# FUTURE WORK
add_section_header(slide, col3_left, content_top + Inches(19), col3_width, "Future Work")

future_text = """• Experimental validation (SPR, ITC binding assays)
• HDX-MS to probe conformational dynamics
• Test additional mutations (G908R, L1007fs)
• Cell-based NOD2 activation assays"""

add_text_box(slide, col3_left + Inches(0.2), content_top + Inches(19.6), col3_width - Inches(0.4), Inches(3),
             future_text, 20, color=WHITE)

# KEY REFERENCES
add_section_header(slide, col3_left, content_top + Inches(23), col3_width, "Key References")

refs_text = """• Hugot et al. Nature 2001 (NOD2-Crohn's link)
• AlphaFold2: Jumper et al. Nature 2021
• FEP methodology: Wang et al. JACS 2015
• All figures generated by author using PyMOL, PLIP"""

add_text_box(slide, col3_left + Inches(0.2), content_top + Inches(23.6), col3_width - Inches(0.4), Inches(2.5),
             refs_text, 18, color=WHITE)

# ============================================================
# BOTTOM STRIP - Acknowledgments
# ============================================================
add_rectangle(slide, Inches(0), Inches(34.5), Inches(48), Inches(1.5), HEADER_BLUE)
add_text_box(slide, Inches(1), Inches(34.7), Inches(46), Inches(1),
             "Computational resources provided by local workstation. All code and analysis performed independently by the author. Contact: [email]",
             18, color=WHITE, align=PP_ALIGN.CENTER)

# Save
output_path = r"C:\Users\vasud\nod2-screening-data\isef_figures\v5_winner\ISEF_2026_RISHAB_STYLE.pptx"
prs.save(output_path)
print(f"SAVED: {output_path}")
print("\nNOW OPEN IN POWERPOINT AND:")
print("1. Insert your actual figures where placeholders are")
print("2. Adjust text sizing if needed")
print("3. The layout matches Rishab Jain's winning style!")
