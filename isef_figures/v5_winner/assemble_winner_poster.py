"""
ISEF 2026 WINNING POSTER - NOD2 Project
Author: Tanmay Vasudeva
Texas Virtual Academy at Hallsville

Based on Rishab Jain's ISEF Grand Award winning poster format.
Dense, information-rich, professional scientific poster.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import nsmap
import os

# Output
OUTPUT_DIR = "C:/Users/vasud/nod2-screening-data/isef_figures/v5_winner/"
FIGURES_DIR = "C:/Users/vasud/nod2-screening-data/isef_figures/v4_isef_style/"
OUTPUT = OUTPUT_DIR + "ISEF_2026_NOD2_WINNER.pptx"

# Colors - Purple/Blue theme like Rishab's poster
HEADER_BG = RGBColor(26, 54, 93)       # Dark navy for headers
SECTION_BG = RGBColor(200, 210, 230)   # Light blue for section boxes
ACCENT = RGBColor(70, 130, 180)        # Steel blue
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
DARK_TEXT = RGBColor(30, 30, 30)
LIGHT_BG = RGBColor(245, 248, 252)     # Very light blue


def add_section_box(slide, title, x, y, width, height):
    """Add a section box with header like Rishab's poster."""
    # Header bar
    header_height = Inches(0.55)
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, width, header_height)
    header.fill.solid()
    header.fill.fore_color.rgb = HEADER_BG
    header.line.fill.background()

    # Header text
    header_text = slide.shapes.add_textbox(x, y + Inches(0.05), width, header_height)
    tf = header_text.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Content box
    content_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                          x, y + header_height,
                                          width, height - header_height)
    content_box.fill.solid()
    content_box.fill.fore_color.rgb = LIGHT_BG
    content_box.line.color.rgb = HEADER_BG
    content_box.line.width = Pt(1.5)

    return y + header_height + Inches(0.1)  # Return content start y


def add_text(slide, text, x, y, width, height, size=14, bold=False, bullet=False):
    """Add text box with optional bullets."""
    text_box = slide.shapes.add_textbox(x, y, width, height)
    tf = text_box.text_frame
    tf.word_wrap = True

    lines = text.strip().split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(6)

    return text_box


def add_figure_with_label(slide, img_path, label, x, y, width, height=None):
    """Add figure with label like 'Fig. 1' """
    if os.path.exists(img_path):
        if height:
            pic = slide.shapes.add_picture(img_path, x, y, width=width, height=height)
        else:
            pic = slide.shapes.add_picture(img_path, x, y, width=width)

        # Add label below
        actual_height = pic.height
        label_box = slide.shapes.add_textbox(x, y + actual_height, width, Inches(0.25))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(9)
        p.font.italic = True
        p.font.color.rgb = DARK_TEXT
        p.alignment = PP_ALIGN.CENTER

        return y + actual_height + Inches(0.3)
    return y


def create_winner_poster():
    """Create ISEF winning-style poster."""

    os.makedirs(OUTPUT_DIR, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(48)  # Landscape 48x36
    prs.slide_height = Inches(36)

    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # ==================== HEADER BANNER ====================
    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(0), Inches(0),
                                     Inches(48), Inches(3.2))
    banner.fill.solid()
    banner.fill.fore_color.rgb = HEADER_BG
    banner.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(46), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Deep Learning Guided Discovery and FEP Validation of NOD2 Binders"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(1.3), Inches(46), Inches(0.6))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "for the R702W Crohn's Variant"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(200, 220, 255)
    p.alignment = PP_ALIGN.CENTER

    # Author
    author_box = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(46), Inches(0.5))
    tf = author_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Tanmay Vasudeva"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # School
    school_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(46), Inches(0.4))
    tf = school_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Texas Virtual Academy at Hallsville"
    p.font.size = Pt(20)
    p.font.italic = True
    p.font.color.rgb = RGBColor(180, 200, 230)
    p.alignment = PP_ALIGN.CENTER

    # One-sentence story box
    story_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Inches(2), Inches(3.4), Inches(44), Inches(0.9))
    story_box.fill.solid()
    story_box.fill.fore_color.rgb = RGBColor(230, 235, 245)
    story_box.line.color.rgb = HEADER_BG
    story_box.line.width = Pt(2)

    story_text = slide.shapes.add_textbox(Inches(2.2), Inches(3.5), Inches(43.6), Inches(0.8))
    tf = story_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "I built a screening pipeline for NOD2 binders and used FEP to quantify how the R702W variant changes binding for specific ligands, with apo MD providing supporting context."
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = DARK_TEXT
    p.alignment = PP_ALIGN.CENTER

    # ==================== COLUMN 1: INTRODUCTION & METHODS ====================
    col1_x = Inches(0.4)
    col_width = Inches(15.2)
    row_y = Inches(4.5)

    # INTRODUCTION
    content_y = add_section_box(slide, "INTRODUCTION", col1_x, row_y, col_width, Inches(4.5))

    intro_text = """CROHN'S DISEASE & NOD2
NOD2 (Nucleotide-binding Oligomerization Domain 2) mutations cause 15-25% of Crohn's disease cases. The R702W variant is the most common mutation, affecting ~10% of patients.

PROBLEM STATEMENT
Current treatments target symptoms, not the underlying molecular dysfunction. No targeted small-molecule modulators exist for NOD2.

HYPOTHESIS
If computational screening identifies NOD2-LRR binders, then FEP can quantify how the R702W mutation affects binding affinity for specific ligands."""
    add_text(slide, intro_text, col1_x + Inches(0.15), content_y, col_width - Inches(0.3), Inches(3.8), size=16)

    # COMPUTATIONAL METHODOLOGY
    method_y = row_y + Inches(4.7)
    content_y = add_section_box(slide, "COMPUTATIONAL METHODOLOGY", col1_x, method_y, col_width, Inches(7))

    # Pipeline figure
    pipeline_img = FIGURES_DIR + "fig_pipeline.png"
    fig_y = add_figure_with_label(slide, pipeline_img, "Fig. 1. Computational screening pipeline overview",
                                   col1_x + Inches(0.3), content_y + Inches(0.1), Inches(14.5))

    method_text = """SCREENING PIPELINE
• 9,566 compounds screened (FDA-approved + natural products)
• GNINA CNN-based docking on AlphaFold2 NOD2 LRR domain
• NOD2-Scout ML model: AUC 0.85-0.93 (scaffold-split cross validation)

MOLECULAR DYNAMICS PROTOCOL
• Engine: OpenMM with AMBER14 + OpenFF 2.1.0 force fields
• Temperature: 310 K | Solvent: TIP3P-FB water model
• Production: 20 ns per replicate | Total: 520 ns across all systems

FREE ENERGY PERTURBATION (FEP)
• Gold standard for binding free energy calculations
• 120 total FEP windows (80 complex + 40 solvent)
• Thermodynamic cycle: WT and R702W systems compared"""
    add_text(slide, method_text, col1_x + Inches(0.15), fig_y + Inches(0.2), col_width - Inches(0.3), Inches(4), size=15)

    # ML PERFORMANCE
    ml_y = method_y + Inches(7.2)
    content_y = add_section_box(slide, "ML MODEL VALIDATION", col1_x, ml_y, col_width, Inches(4.8))

    ml_img = FIGURES_DIR + "fig_ml_performance.png"
    add_figure_with_label(slide, ml_img, "Fig. 2. NOD2-Scout model performance across 5-fold scaffold-split CV",
                          col1_x + Inches(0.5), content_y + Inches(0.1), Inches(14))

    # ==================== COLUMN 2: KEY RESULTS ====================
    col2_x = Inches(16)

    # FEP RESULTS (KEY FINDING)
    content_y = add_section_box(slide, "FEP BINDING FREE ENERGY (KEY RESULT)", col2_x, row_y, col_width, Inches(8))

    fep_img = FIGURES_DIR + "fig_fep.png"
    fig_y = add_figure_with_label(slide, fep_img, "Fig. 3. FEP binding free energy difference (R702W - WT)",
                                   col2_x + Inches(0.5), content_y + Inches(0.1), Inches(10))

    # FEP Results Table
    fep_table_text = """FEP RESULTS SUMMARY (Gold Standard Thermodynamics)

COMPOUND        | ΔG_WT (kcal/mol) | ΔG_R702W      | ΔΔG        | Significance
----------------|------------------|---------------|------------|-------------
Febuxostat      | -10.36 ± 0.18    | -8.02 ± 0.19  | +2.34      | *** (8σ, p<0.001)
Bufadienolide   | -15.22 ± 0.26    | -15.66 ± 0.26 | -0.44      | n.s. (1.8σ)

KEY FINDING: Febuxostat shows 50× weaker binding in R702W mutant
             Bufadienolide shows no significant binding difference"""
    add_text(slide, fep_table_text, col2_x + Inches(0.15), fig_y + Inches(0.3), col_width - Inches(0.3), Inches(3.5), size=15)

    # MD BINDING VALIDATION
    md_y = row_y + Inches(8.2)
    content_y = add_section_box(slide, "MD BINDING VALIDATION", col2_x, md_y, col_width, Inches(5))

    pocket_img = FIGURES_DIR + "fig_pocket_occupancy.png"
    fig_y = add_figure_with_label(slide, pocket_img, "Fig. 4. Pocket occupancy during 20 ns MD simulations",
                                   col2_x + Inches(0.5), content_y + Inches(0.1), Inches(10))

    md_text = """POCKET OCCUPANCY (Ligand within 5Å of binding site)
• Febuxostat: 70% mean occupancy | H-bonds: 2.02 ± 1.78
• Bufadienolide: 80% mean occupancy | H-bonds: 1.99 ± 1.74
• Decoy (control): 0% occupancy - validates binding specificity"""
    add_text(slide, md_text, col2_x + Inches(0.15), fig_y + Inches(0.1), col_width - Inches(0.3), Inches(1.5), size=15)

    # KEY BINDING RESIDUES
    res_y = md_y + Inches(5.2)
    content_y = add_section_box(slide, "KEY BINDING RESIDUES", col2_x, res_y, col_width, Inches(3.3))

    residue_text = """IDENTIFIED FROM CONTACT ANALYSIS (Phase 7 MD)

Residue    | Contact Frequency | Role
-----------|-------------------|---------------------------
GLU1008    | 50-71%           | H-bond acceptor
ARG1037    | 71-86%           | Electrostatic anchor
ASP1011    | 67-93%           | H-bond acceptor
ARG1034    | 82-88%           | Secondary contact
LEU1007    | 61-65%           | Hydrophobic packing"""
    add_text(slide, residue_text, col2_x + Inches(0.15), content_y + Inches(0.1), col_width - Inches(0.3), Inches(2.8), size=15)

    # ==================== COLUMN 3: DYNAMICS & CONCLUSIONS ====================
    col3_x = Inches(31.6)

    # APO MD DYNAMICS
    content_y = add_section_box(slide, "APO MD DYNAMICS (SUPPORTING CONTEXT)", col3_x, row_y, col_width, Inches(6.5))

    dynamics_img = FIGURES_DIR + "fig_dynamics.png"
    fig_y = add_figure_with_label(slide, dynamics_img, "Fig. 5. Apo MD comparison: WT vs R702W backbone dynamics",
                                   col3_x + Inches(0.3), content_y + Inches(0.1), Inches(14.5))

    dynamics_text = """APO SIMULATION RESULTS (No Ligand Bound)

Metric              | WT apo (n=4) | R702W apo (n=3)
--------------------|--------------|----------------
Mean RMSD           | 4.05 Å       | 3.27 Å
% Frames > 5Å       | 12.6%        | 0%

INTERPRETATION: R702W shows reduced conformational sampling
compared to WT, consistent with mutation-sensitive binding in FEP."""
    add_text(slide, dynamics_text, col3_x + Inches(0.15), fig_y + Inches(0.1), col_width - Inches(0.3), Inches(2.2), size=15)

    # RMSD TRAJECTORIES
    trace_y = row_y + Inches(6.7)
    content_y = add_section_box(slide, "REPRESENTATIVE MD TRAJECTORIES", col3_x, trace_y, col_width, Inches(4.3))

    trace_img = FIGURES_DIR + "fig_rmsd_trace.png"
    add_figure_with_label(slide, trace_img, "Fig. 6. RMSD time series showing WT crosses 5Å threshold, R702W does not",
                          col3_x + Inches(0.3), content_y + Inches(0.1), Inches(14.5))

    # ==================== BOTTOM ROW ====================
    bottom_y = Inches(28.5)
    bottom_height = Inches(7)

    # CONCLUSIONS
    content_y = add_section_box(slide, "CONCLUSIONS", col1_x, bottom_y, Inches(17), bottom_height)

    conclusions_text = """SCIENTIFIC CONTRIBUTIONS

1. PIPELINE DEVELOPMENT
   Built deep learning screening pipeline for 9,566 compounds against NOD2 LRR domain

2. FEP QUANTIFIES MUTATION-SENSITIVE BINDING (Primary Finding)
   • Febuxostat: ΔΔG = +2.34 kcal/mol → 50× weaker in R702W (p<0.001)
   • Bufadienolide: ΔΔG = -0.44 kcal/mol → No significant change
   • Demonstrates ligand-dependent binding differences

3. APO DYNAMICS PROVIDE MECHANISTIC CONTEXT (Supporting Finding)
   • WT: 12.6% frames exceed 5Å RMSD threshold
   • R702W: 0% frames exceed threshold
   • Reduced conformational sampling in mutant

4. STRUCTURAL INSIGHTS
   Key binding residues identified: GLU1008, ARG1037, ASP1011"""
    add_text(slide, conclusions_text, col1_x + Inches(0.15), content_y + Inches(0.1), Inches(16.5), Inches(6.5), size=16)

    # LIMITATIONS
    lim_x = Inches(17.8)
    content_y = add_section_box(slide, "LIMITATIONS", lim_x, bottom_y, Inches(9.5), Inches(3.3))

    limitations_text = """• FEP performed on 2 representative ligands only
• Apo MD simulations (no ligand-bound comparisons)
• Computational predictions require experimental validation
• Force field and sampling limitations inherent to MD"""
    add_text(slide, limitations_text, lim_x + Inches(0.15), content_y + Inches(0.1), Inches(9), Inches(2.8), size=15)

    # FUTURE WORK
    future_y = bottom_y + Inches(3.5)
    content_y = add_section_box(slide, "FUTURE WORK", lim_x, future_y, Inches(9.5), Inches(3.3))

    future_text = """• SPR/ITC experimental binding validation
• Test additional NOD2 mutations (G908R, L1007fs)
• HDX-MS conformational studies
• Analog optimization for lead compounds"""
    add_text(slide, future_text, lim_x + Inches(0.15), content_y + Inches(0.1), Inches(9), Inches(2.8), size=15)

    # KEY REFERENCES
    ref_x = Inches(27.7)
    content_y = add_section_box(slide, "KEY REFERENCES", ref_x, bottom_y, Inches(19.5), Inches(3.3))

    refs_text = """[1] Hugot JP, et al. Association of NOD2 leucine-rich repeat variants with susceptibility to Crohn's disease. Nature. 2001.
[2] Jumper J, et al. Highly accurate protein structure prediction with AlphaFold. Nature. 2021.
[3] McNutt AT, et al. GNINA 1.0: molecular docking with deep learning. J Cheminform. 2021.
[4] Shirts MR, et al. Best practices for alchemical free energy calculations. Living J Comp Mol Sci. 2019."""
    add_text(slide, refs_text, ref_x + Inches(0.15), content_y + Inches(0.1), Inches(19), Inches(2.8), size=14)

    # ACKNOWLEDGMENTS
    ack_y = bottom_y + Inches(3.5)
    content_y = add_section_box(slide, "ACKNOWLEDGMENTS", ref_x, ack_y, Inches(19.5), Inches(3.3))

    ack_text = """• Computational resources provided by vast.ai GPU cloud
• AlphaFold2 structure from DeepMind/EMBL-EBI
• OpenMM and GROMACS development teams
• Open-source cheminformatics community (RDKit, OpenFF)"""
    add_text(slide, ack_text, ref_x + Inches(0.15), content_y + Inches(0.1), Inches(19), Inches(2.8), size=15)

    # Save
    prs.save(OUTPUT)
    print(f"WINNER POSTER saved to: {OUTPUT}")


if __name__ == "__main__":
    create_winner_poster()
    print("\nPoster created in ISEF winning style!")
    print("Open in PowerPoint to add PyMOL 3D structure images.")
